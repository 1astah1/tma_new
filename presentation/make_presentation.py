import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn

PRES_DIR = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS_DIR = os.path.join(PRES_DIR, "screenshots")
ASSETS_DIR = r"D:\TMA_Seill\tma_new\CoinMint"

GOLD = RGBColor(0xFF, 0xD7, 0x00)
BLACK = RGBColor(0x0D, 0x0D, 0x0D)
DARK_GRAY = RGBColor(0x1A, 0x1A, 0x1A)
MID_GRAY = RGBColor(0x2A, 0x2A, 0x2A)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_CARD = RGBColor(0x15, 0x15, 0x15)

def add_transition(slide, kind="fade", speed="medium"):
    t = slide._element.makeelement(qn("p:transition"), {})
    t.set("spd", speed)
    t.set("advClick", "1")
    child = etree.SubElement(t, qn(f"p:{kind}"))
    slide._element.append(t)

def add_animation(slide, shape_id, anim_type="fade", delay_ms=0):
    timing = slide._element.find(qn("p:timing"))
    if timing is None:
        timing = etree.SubElement(slide._element, qn("p:timing"))
    # Build the animation XML
    tnLst = timing.find(qn("p:tnLst"))
    if tnLst is None:
        tnLst = etree.SubElement(timing, qn("p:tnLst"))
    # We'll keep it simple - just add a basic animation reference
    # python-pptx doesn't support this well, so we skip for now

def find_logo():
    for root, dirs, files in os.walk(ASSETS_DIR):
        for f in files:
            low = f.lower()
            if "1 без" in low or "1 " in low[:5] and "фон" in low:
                return os.path.join(root, f)
    return None

def find_banner():
    for root, dirs, files in os.walk(ASSETS_DIR):
        for f in files:
            low = f.lower()
            if "баннер" in low and "с назв" in low and f.endswith(".png"):
                return os.path.join(root, f)
    return None

def find_product_images():
    products = {}
    for root, dirs, files in os.walk(ASSETS_DIR):
        for f in files:
            if f.endswith(".png") and "100x100" not in f and len(f) < 15:
                name = os.path.splitext(f)[0]
                products[name] = os.path.join(root, f)
    return products

def add_bg(slide, color=None):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or DARK_GRAY

def add_title_bar(slide, title, y=0, height=0.85):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y), Emu(10080000), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_GRAY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.color.rgb = GOLD
    p.font.bold = True
    return shape

def add_accent_line(slide, left, top, width=2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape

def add_bullet_text(slide, bullets, left=0.6, top=1.1, width=8.8, height=6, font_size=13):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(b, tuple):
            pp.text = b[0]
            pp.font.size = Pt(font_size + 1)
            pp.font.color.rgb = b[1] if len(b) > 1 else WHITE
            pp.font.bold = True
        else:
            pp.text = b
            pp.font.size = Pt(font_size)
            pp.font.color.rgb = WHITE
        pp.space_after = Pt(5)
    return txBox

def add_image_safe(slide, path, left, top, width=None, height=None):
    if not path or not os.path.exists(path):
        return None
    try:
        kw = {}
        if width:
            kw["width"] = Inches(width)
        if height:
            kw["height"] = Inches(height)
        return slide.shapes.add_picture(path, Inches(left), Inches(top), **kw)
    except:
        return None

def make_slide_title(prs, main_text, sub_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BLACK)
    add_transition(slide, "fade")
    logo = find_logo()
    if logo and os.path.exists(logo):
        try:
            slide.shapes.add_picture(logo, 0, 0, Emu(10080000))
        except:
            pass
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(10080000), Inches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = BLACK
    overlay.line.fill.background()
    # Make it semi-transparent by adjusting alpha
    solidFill = overlay._element.find(qn("a:solidFill"))
    if solidFill is not None:
        clr = etree.SubElement(solidFill, qn("a:srgbClr"))
        clr.set("val", "0D0D0D")
        alpha = etree.SubElement(clr, qn("a:alpha"))
        alpha.set("val", "60000")

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = main_text
    p.font.size = Pt(44)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if sub_text:
        add_accent_line(slide, 4, 3.6, 2)
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.9), Inches(8), Inches(1))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = sub_text
        p2.font.size = Pt(18)
        p2.font.color.rgb = LIGHT_GRAY
        p2.alignment = PP_ALIGN.CENTER
    return slide

def make_slide_content(prs, title, bullets, screenshot_path=None, img_side="right"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_GRAY)
    add_transition(slide, "push" if img_side == "right" else "fade")
    add_title_bar(slide, title)

    if screenshot_path and os.path.exists(screenshot_path):
        if img_side == "right":
            add_image_safe(slide, screenshot_path, 5.0, 1.2, width=4.7)
            add_bullet_text(slide, bullets, left=0.6, top=1.2, width=4.2, font_size=12)
        else:
            add_image_safe(slide, screenshot_path, 0.3, 1.2, height=5.5)
            add_bullet_text(slide, bullets, left=5.3, top=1.2, width=4.5, font_size=12)
    else:
        add_bullet_text(slide, bullets, left=0.6, top=1.2, width=8.8, font_size=13)
    return slide

def make_slide_products(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_GRAY)
    add_transition(slide, "fade")
    add_title_bar(slide, "Каталог товаров")

    products = find_product_images()
    all_products = [("BC.png", "Battlefield"), ("CP.png", "Cyberpunk"), ("WZ100x100.png", "Warzone"),
                    ("PS Store", "PS Store"), ("Xbox", "Xbox Store"), ("EA PLAY", "EA Play")]

    x, y = 0.3, 1.3
    cols, rows, spacing_x, spacing_y = 4, 2, 2.4, 3.0
    placed = 0
    for root, dirs, files in os.walk(ASSETS_DIR):
        for f in sorted(files):
            if placed >= 8:
                break
            if "100x100" not in f and f.endswith(".png") and os.path.getsize(os.path.join(root, f)) > 50000:
                col = placed % cols
                row = placed // cols
                px = x + col * spacing_x
                py = y + row * spacing_y
                try:
                    slide.shapes.add_picture(os.path.join(root, f), Inches(px), Inches(py), Inches(2.0), Inches(2.0))
                except:
                    pass
                placed += 1

    bullets = [
        "Сетка товаров с изображениями",
        "Фильтр по платформе (PS4/PS5/Xbox/PC)",
        "Фильтр по типу (игры/подписки/валюта)",
        "Детальная карточка с описанием",
        "Выбор доставки: ключ или активация",
        "Промокоды на скидку",
    ]
    add_bullet_text(slide, bullets, left=0.5, top=5.3, width=9, font_size=11)
    return slide

def make_slide_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_GRAY)
    add_transition(slide, "fade")
    add_title_bar(slide, "Сравнение: TMA vs Telegram Bot")

    headers = ["Критерий", "@shopcoinmint", "TMA Shop"]
    header_colors = [LIGHT_GRAY, LIGHT_GRAY, GOLD]
    rows = [
        ["Интерфейс", "Текстовые кнопки", "Нативный UI (React)"],
        ["Корзина", "Нет", "Полноценная корзина"],
        ["Каталог", "Список текстом", "Сетка с картинками"],
        ["Оплата", "Только карта", "СБП / Карта / Крипта"],
        ["Чат поддержки", "Нет", "Встроенный чат"],
        ["Статус заказа", "Текст", "Прогресс-бар + степпер"],
        ["Личный кабинет", "Нет", "История заказов"],
        ["Админ-панель", "Команды в чат", "React Admin"],
        ["Выдача ключей", "Вручную", "Автоматически"],
        ["Шифрование", "Нет", "AES-256"],
        ["Промокоды", "Нет", "Гибкая система"],
        ["Масштабируемость", "Низкая", "Высокая (Go+PG)"],
    ]

    table_left = 0.7
    table_top = 1.2
    col_widths = [2.6, 3.2, 3.2]
    row_height = 0.4

    x = table_left
    for j, h in enumerate(headers):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(table_top), Inches(col_widths[j]), Inches(row_height))
        s.fill.solid()
        s.fill.fore_color.rgb = MID_GRAY
        s.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        tf = s.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = h
        p.font.size = Pt(11)
        p.font.color.rgb = header_colors[j]
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        x += col_widths[j]

    for i, row in enumerate(rows):
        y = table_top + row_height + i * row_height
        bg = RGBColor(0x18, 0x18, 0x18) if i % 2 == 0 else RGBColor(0x12, 0x12, 0x12)
        x = table_left
        for j, cell in enumerate(row):
            s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(col_widths[j]), Inches(row_height))
            s.fill.solid()
            s.fill.fore_color.rgb = bg
            s.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
            tf = s.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(10)
            p.font.color.rgb = GOLD if cell == "TMA Shop" or j == 2 else WHITE
            p.font.bold = (cell == "TMA Shop")
            p.alignment = PP_ALIGN.CENTER
            x += col_widths[j]

    return slide

def make_slide_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_GRAY)
    add_transition(slide, "fade")
    add_title_bar(slide, "Архитектура")

    arch = [
        ("TMA Frontend", "React + Vite", ":5173", "Интерфейс покупателя внутри Telegram"),
        ("Admin Panel", "React Admin\n+ MUI", ":5174", "Панель управления"),
        ("Backend API", "Go + Chi\n+ sqlx", ":8081", "REST API, шифрование, логика"),
        ("PostgreSQL", "Database\n16", ":5432", "Заказы, товары, ключи, юзеры"),
    ]

    y = 1.3
    for name, tech, port, desc in arch:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(y), Inches(2.5), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = MID_GRAY
        box.line.color.rgb = GOLD
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{name}"
        p.font.size = Pt(13)
        p.font.color.rgb = GOLD
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = tech
        p2.font.size = Pt(9)
        p2.font.color.rgb = LIGHT_GRAY
        p2.alignment = PP_ALIGN.CENTER

        dbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(y), Inches(6.0), Inches(1.2))
        dbox.fill.solid()
        dbox.fill.fore_color.rgb = BLACK
        dbox.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        dtf = dbox.text_frame
        dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(12)
        dp.font.color.rgb = WHITE
        dp2 = dtf.add_paragraph()
        dp2.text = f"Порт: {port}"
        dp2.font.size = Pt(10)
        dp2.font.color.rgb = LIGHT_GRAY

        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.9), Inches(y + 0.4), Inches(0.5), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

        y += 1.5

    return slide

def make_slide_business(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_GRAY)
    add_transition(slide, "fade")
    add_title_bar(slide, "Бизнес-ценность")

    items = [
        ("\u2191 Конверсия", "Удобный UI повышает % завершения покупок"),
        ("\u2193 Нагрузка на админов", "Автоматическая выдача ключей"),
        ("\u2191 Лояльность", "Чат поддержки, прозрачный статус"),
        ("\u2191 Средний чек", "Корзина — несколько товаров сразу"),
        ("\u2191 Безопасность", "AES-256, защита от брутфорса"),
        ("\u2193 Ошибки", "Автоматизация исключает человеческий фактор"),
        ("\u2191 Масштабируемость", "Go + PG выдерживают высокие нагрузки"),
        ("\u2191 Аналитика", "Дашборд, статистика, аудит"),
    ]

    y = 1.2
    for i, (title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 4.8
        yy = y + row * 1.5

        num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(yy), Inches(0.4), Inches(0.4))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = GOLD
        num_box.line.fill.background()
        ntf = num_box.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = ntf.paragraphs[0]
        np.text = str(i + 1)
        np.font.size = Pt(11)
        np.font.color.rgb = BLACK
        np.font.bold = True
        np.alignment = PP_ALIGN.CENTER

        txBox = slide.shapes.add_textbox(Inches(x + 0.55), Inches(yy), Inches(4.1), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.color.rgb = GOLD
        p.font.bold = True
        p.space_after = Pt(2)
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = LIGHT_GRAY

    return slide

def create_presentation():
    print("=" * 60)
    print("COIN MINT TMA SHOP — генерация презы")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width = Emu(10080000)
    prs.slide_height = Emu(7560000)

    shots = SCREENSHOTS_DIR

    # 1 — Title
    make_slide_title(prs, "COIN MINT", "Telegram Mini App — платформа для продажи цифровых товаров")

    # 2 — Проблема
    make_slide_content(prs, "Проблема: @shopcoinmint Bot", [
        "Текстовый интерфейс — неинтуитивно и медленно",
        "Нет корзины — только по одному товару за раз",
        "Нет визуального статуса заказа",
        "Нет чата с поддержкой",
        "Администрирование через команды в Telegram",
        "Ручная выдача ключей — ошибки и задержки",
        "Нет личного кабинета и истории покупок",
        "Сложно масштабировать",
    ])

    # 3 — Решение
    make_slide_content(prs, "Решение: Telegram Mini App (TMA)", [
        "Нативный UI внутри Telegram — React + TypeScript",
        "Полноценная корзина с промокодами",
        "Визуальный прогресс заказа (степпер)",
        "Встроенный чат с поддержкой",
        "React Admin панель с дашбордом",
        "Автоматическая выдача ключей",
        "AES-256 шифрование данных",
        "Go + PostgreSQL — высокая производительность",
    ], os.path.join(shots, "tma_products.png"))

    # 4 — Каталог
    make_slide_products(prs)

    # 5 — Оформление
    make_slide_content(prs, "Оформление и оплата", [
        "Корзина с несколькими товарами",
        "Способы оплаты: СБП, Карта, Криптовалюта",
        "Отображение реквизитов для оплаты",
        "Загрузка чека (фото или PDF, до 5MB)",
        "Автоматическая проверка оплаты",
        "Поддержка batch-заказов",
    ], os.path.join(shots, "tma_checkout.png"))

    # 6 — Статус заказа
    make_slide_content(prs, "Статус заказа и отслеживание", [
        "Progress bar с процентом выполнения",
        "Пошаговый степпер для key и activation",
        "17 статусов с валидными переходами",
        "Ввод логина/пароля и 2FA для активации",
        "Получение ключа активации (копирование)",
        "Автоуведомление о смене статуса",
    ], os.path.join(shots, "tma_order_status.png"))

    # 7 — Чат
    make_slide_content(prs, "Чат поддержки", [
        "Встроенный чат внутри каждого заказа",
        "Автоматическая отправка ключа в чат",
        "Шаблоны сообщений для админа",
        "Polling обновлений (3 сек)",
        "Особое отображение ключей активации",
    ], os.path.join(shots, "admin_order_detail.png"), img_side="left")

    # 8 — Дашборд
    make_slide_content(prs, "Админ-панель: Дашборд", [
        "Общая статистика по заказам",
        "Количество заказов по статусам",
        "Выручка за период",
        "Последние заказы",
        "Графики и метрики",
    ], os.path.join(shots, "admin_dashboard.png"))

    # 9 — Управление заказами
    make_slide_content(prs, "Админ-панель: Заказы", [
        "Список с фильтрацией по статусу",
        "Детальный просмотр заказа",
        "Кнопки действий по статусу",
        "История изменений статуса",
        "Чат с клиентом",
        "Загрузка и просмотр чеков",
        "Выдача ключей (выбор из доступных)",
        "Дешифровка credentials",
    ], os.path.join(shots, "admin_orders_list.png"))

    # 10 — Детали заказа
    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s10, DARK_GRAY)
    add_transition(s10, "fade")
    add_title_bar(s10, "Детали заказа и статус-машина")

    img1 = os.path.join(shots, "admin_order_detail.png")
    img2 = os.path.join(shots, "admin_activation_order.png")
    if os.path.exists(img1):
        add_image_safe(s10, img1, 0.2, 1.1, width=3.0)
    if os.path.exists(img2):
        add_image_safe(s10, img2, 3.4, 1.1, width=3.0)
    add_bullet_text(s10, [
        "17 статусов с валидными переходами",
        "Действия зависят от статуса",
        "Информация о клиенте и товаре",
        "Отображение чека об оплате",
        "Данные аккаунта (логин/пароль/2FA)",
        "Назначенный администратор",
        "История всех изменений",
    ], left=6.8, top=1.1, width=3.0, font_size=11)

    # 11 — Товары и ключи
    s11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s11, DARK_GRAY)
    add_transition(s11, "fade")
    add_title_bar(s11, "Товары и ключи")
    img_p = os.path.join(shots, "admin_products.png")
    img_k = os.path.join(shots, "admin_keys.png")
    if os.path.exists(img_p):
        add_image_safe(s11, img_p, 0.2, 1.1, width=3.0)
    if os.path.exists(img_k):
        add_image_safe(s11, img_k, 3.4, 1.1, width=3.0)
    add_bullet_text(s11, [
        "CRUD товаров с изображениями",
        "Множественные способы доставки",
        "Варианты товара (разные цены)",
        "Импорт ключей пачкой",
        "Статусы ключей (available/sold)",
        "Освобождение ключей при отмене",
    ], left=6.8, top=1.1, width=3.0, font_size=11)

    # 12 — Пользователи
    make_slide_content(prs, "Пользователи и безопасность", [
        "Список пользователей с поиском",
        "История заказов пользователя",
        "Блокировка / бан пользователей",
        "Админ-заметки к пользователям",
        "Система ролей администраторов",
        "Rate limiting, brute-force защита",
        "Аудит всех действий админов",
    ], os.path.join(shots, "admin_users.png"))

    # 13 — Архитектура
    make_slide_arch(prs)

    # 14 — Сравнение
    make_slide_comparison(prs)

    # 15 — Бизнес-ценность
    make_slide_business(prs)

    # 16 — Спасибо
    s16 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s16, BLACK)
    add_transition(s16, "fade")

    txBox = s16.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Спасибо!"
    p.font.size = Pt(44)
    p.font.color.rgb = GOLD
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_accent_line(s16, 4, 4.1, 2)

    txBox2 = s16.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Готовы ответить на ваши вопросы"
    p2.font.size = Pt(18)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    output_path = os.path.join(PRES_DIR, "COIN_MINT_TMA_Presentation.pptx")
    prs.save(output_path)
    print(f"\nСохранено: {output_path}")
    print("=" * 60)

if __name__ == "__main__":
    create_presentation()
